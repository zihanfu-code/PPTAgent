import re
from dataclasses import dataclass
from enum import Enum, auto
from types import MappingProxyType
from typing import Callable, ClassVar, Optional, Union

from lxml import etree
from pptx.dml.fill import FillFormat
from pptx.dml.line import LineFormat
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.parts.slide import SlidePart
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape as PPTXGroupShape
from pptx.shapes.picture import Picture as PPTXPicture
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
from pptx.slide import Slide as PPTXSlide
from pptx.slide import _Background
from pptx.text.text import _Paragraph
from pptx.util import Length

from pptagent.utils import (
    Config,
    dict_to_object,
    package_join,
    parse_groupshape,
    parsing_image,
    pjoin,
    runs_merge,
)

INDENT = "\t"


def shape_normalize(shape: BaseShape):
    """
    This function is used to filter out those malfunctioned attrs.
    """
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.hyperlink.address = None


class ClosureType(Enum):
    CLONE = auto()
    REPLACE = auto()
    DELETE = auto()
    STYLE = auto()
    MERGE = auto()

    def __str__(self):
        return self.name.lower()

    @classmethod
    def to_default_dict(cls):
        return {key: [] for key in cls}


@dataclass
class StyleArg:
    """
    A class to represent style arguments for HTML conversion.
    """

    paragraph_id: bool = True
    element_id: bool = True
    font_style: bool = True
    fill_style: bool = True
    area: bool = False
    size: bool = False
    geometry: bool = False
    show_name: bool = False
    show_image: bool = True
    show_empty: bool = False
    show_content: bool = True
    show_semantic_name: bool = False

    @classmethod
    def all_true(cls) -> "StyleArg":
        """
        Create a StyleArg instance with all options enabled.

        Returns:
            StyleArg: A StyleArg instance with all options enabled.
        """
        return cls(
            area=True,
            size=True,
            geometry=True,
            show_semantic_name=True,
        )


class Fill:
    """
    A class to represent a fill.
    """

    def __init__(
        self,
        fill_type: MSO_FILL_TYPE,
        fill_str: str,
        fill_xml: str,
        image_path: Optional[str] = None,
    ):
        self.fill_type = fill_type
        self.fill_str = fill_str
        self.fill_xml = fill_xml
        self.image_path = image_path

    @classmethod
    def from_shape(cls, fill: Optional[FillFormat], part: SlidePart, config: Config):
        if fill is None or fill.type is None or fill.type == MSO_FILL_TYPE.BACKGROUND:
            return cls(MSO_FILL_TYPE.BACKGROUND, "", None)

        fill_str = "Fill: " + str(fill.value)
        fill_xml = fill._xPr.xml
        fill_type = fill.type
        image_path = None
        if fill_type == MSO_FILL_TYPE.PICTURE:
            image = part.get_image(fill._fill.rId)
            image_path = pjoin(config.IMAGE_DIR, f"{image.sha1}.{image.ext}")
            image_path = parsing_image(image, image_path)
        return cls(fill_type, fill_str, fill_xml, image_path)

    # We pass an element with fill attribute instead of a fill object because `python-pptx` automatically creates a fill object when accessing this attribute, which would cause inconsistency
    def build(
        self, fill_ele: LineFormat | _Background | BaseShape, part: SlidePart
    ) -> None:
        """
        Build the fill in a shape.
        Args:
            shape (BaseShape): The shape to apply fill to.
            fill_xml (Optional[str]): The fill XML to apply.
        """
        if self.fill_type == MSO_FILL_TYPE.BACKGROUND:
            return
        fill = fill_ele.fill
        if self.fill_type == MSO_FILL_TYPE.PICTURE:
            fill.blip()
            _, rId = part.get_or_add_image_part(self.image_path)
            fill.rId = rId
        else:
            new_element = etree.fromstring(self.fill_xml)
            fill._xPr.getparent().replace(fill._xPr, new_element)

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the fill to HTML.
        """


class Line:
    """
    A class to represent a line.
    """

    def __init__(self, fill: Fill, line_width: float, line_dash_style: str):
        self.fill = fill
        self.line_width = line_width
        self.line_dash_style = line_dash_style

    @classmethod
    def from_shape(cls, line: Optional[LineFormat], part: SlidePart, config: Config):
        line_fill = getattr(line, "fill", None)
        if line_fill is None:
            return cls(Fill(MSO_FILL_TYPE.BACKGROUND, "", None), 0, "")
        fill = Fill.from_shape(line_fill, part, config)
        line_width = line.width
        line_dash_style = line.dash_style
        return cls(fill, line_width, line_dash_style)

    def build(self, line: LineFormat, part: SlidePart) -> None:
        """
        Build the line in a shape.
        """
        if self.fill.fill_type == MSO_FILL_TYPE.BACKGROUND:
            return
        self.fill.build(line, part)
        line.width = self.line_width
        line.dash_style = self.line_dash_style


class Background(Fill):
    """
    A class to represent a slide background.
    """

    shape_idx: int = -1

    @classmethod
    def from_slide(cls, slide: PPTXSlide, config: Config) -> "Background":
        """
        Build the background in a slide.

        Args:
            slide (PPTXSlide): The slide to build the background in.
        """
        background = slide.background
        return cls.from_shape(background.fill, slide.part, config)

    def build(self, slide: PPTXSlide) -> None:
        """
        Build the background in a slide.
        """
        super().build(slide.background, slide.part)

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the background to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the background.
        """
        raise NotImplementedError("Background to HTML conversion is not implemented")

    @property
    def closures(self) -> list:
        """
        Get the closure for the background.
        """
        return []


@dataclass
class Closure:
    """
    A class to represent a closure that can be applied to a shape.
    """

    closure: Callable[[BaseShape], None]
    paragraph_id: int = -1

    def apply(self, shape: BaseShape) -> None:
        """
        Apply the closure to a shape.

        Args:
            shape (BaseShape): The shape to apply the closure to.
        """
        self.closure(shape)

    def __gt__(self, other: "Closure") -> bool:
        """
        Compare closures based on paragraph_id.

        Args:
            other (Closure): Another closure to compare with.

        Returns:
            bool: True if this closure's paragraph_id is greater than the other's.
        """
        if self.paragraph_id != other.paragraph_id:
            return self.paragraph_id > other.paragraph_id


@dataclass
class Font:
    name: str
    color: str
    size: Length
    bold: bool
    italic: bool
    underline: bool
    strikethrough: bool

    def update(self, other: "Font"):
        """
        Merge a list of fonts into a single font.
        """
        for key, value in other.__dict__.items():
            if getattr(self, key) is None:
                setattr(self, key, value)

    def override(self, other: "Font"):
        """
        Merge a list of fonts into a single font.
        """
        for key, value in other.__dict__.items():
            if value is not None:
                setattr(self, key, value)

    def unify(self, others: list["Font"], clear_others: bool = False):
        """
        Merge a list of fonts into a single font.
        """
        if len(others) == 0:
            return
        for key in list(self.__dict__.keys()):
            values = [d.__dict__[key] for d in others]
            if not all(value == values[0] for value in values):
                continue
            setattr(self, key, values[0])
            if not clear_others:
                continue
            for d in others:
                setattr(d, key, None)

    def to_style(self) -> str:
        """
        Convert a font dictionary to a CSS style string.

        Args:
            font (Dict[str, Any]): The font dictionary.

        Returns:
            str: The CSS style string.
        """
        styles = []
        if self.size:
            styles.append(f"font-size: {self.size}pt")

        if self.color:
            styles.append(f"color: #{self.color}")

        if self.bold:
            styles.append("font-weight: bold")

        if self.italic:
            styles.append("font-style: italic")

        return "; ".join(styles)


class Paragraph:
    """
    A class to represent a paragraph in a text frame.
    """

    def __init__(self, paragraph: _Paragraph, idx: int):
        """
        Initialize a Paragraph.

        Args:
            paragraph (_Paragraph): The paragraph object.
            idx (int): The index of the paragraph.
        """
        run = runs_merge(paragraph)
        self.idx = idx
        self.real_idx = idx
        self.bullet = paragraph.bullet
        if run is None:
            self.idx = -1
            return
        self.font = Font(**paragraph.font.get_attrs())
        self.font.override(Font(**run.font.get_attrs()))
        self.text = re.sub(r"(_x000B_|\\x0b)", " ", paragraph.text)

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the paragraph to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the paragraph.

        Raises:
            ValueError: If the paragraph is not valid.
        """
        if self.idx == -1:
            raise ValueError(f"paragraph {self.idx} is not valid")
        tag = "li" if self.bullet else "p"
        id_str = f" id='{self.idx}'" if style_args.paragraph_id else ""
        font_style = self.font.to_style()
        style_str = (
            f" style='{font_style}'" if style_args.font_style and font_style else ""
        )
        if self.bullet:
            style_str += f" bullet-type='{self.bullet}'"
        return f"<{tag}{id_str}{style_str}>{self.text}</{tag}>"

    def __repr__(self) -> str:
        """
        Get a string representation of the paragraph.

        Returns:
            str: A string representation of the paragraph.
        """
        return f"Paragraph-{self.idx}: {self.text}"


class TextFrame:
    """
    A class to represent a text frame in a shape.
    """

    def __init__(self, shape: BaseShape, level: int):
        """
        Initialize a TextFrame.

        Args:
            shape (BaseShape): The shape containing the text frame.
            level (int): The indentation level.
        """
        if not shape.has_text_frame:
            self.is_textframe = False
            return
        self.paragraphs = [
            Paragraph(paragraph, idx)
            for idx, paragraph in enumerate(shape.text_frame.paragraphs)
        ]
        para_offset = 0
        for para in self.paragraphs:
            if para.idx == -1:
                para_offset += 1
            else:
                para.idx = para.idx - para_offset
        if len(self.paragraphs) == 0:
            self.is_textframe = False
            return
        self.level = level
        self.text = shape.text
        self.is_textframe = True
        self.extents = shape.text_frame._extents
        self.font = Font(**shape.text_frame.font.get_attrs())
        self.font.unify([para.font for para in self.paragraphs if para.idx != -1])

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the text frame to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the text frame.
        """
        if not self.is_textframe:
            return ""
        repr_list = [
            para.to_html(style_args) for para in self.paragraphs if para.idx != -1
        ]
        return "\n".join([INDENT * self.level + repr for repr in repr_list])

    def __repr__(self) -> str:
        """
        Get a string representation of the text frame.

        Returns:
            str: A string representation of the text frame.
        """
        if not self.is_textframe:
            return "TextFrame: null"
        return f"TextFrame: {self.paragraphs}"

    def __len__(self) -> int:
        """
        Get the length of the text in the text frame.

        Returns:
            int: The length of the text.
        """
        if not self.is_textframe:
            return 0
        return len(self.text)


@dataclass
class ShapeElement:
    """
    Base class for shape elements in a presentation.
    """

    config: Config
    slide_idx: int
    shape_idx: int
    style: dict
    data: list
    text_frame: TextFrame
    level: int
    slide_area: float
    xml: str
    fill: Fill
    line: Line
    shape: BaseShape
    _closures: dict[ClosureType, list[Closure]]

    @classmethod
    def from_shape(
        cls: type["ShapeElement"],
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        config: Config,
        slide_area: float,
        shape_cast: dict[MSO_SHAPE_TYPE, type["ShapeElement"] | None],
        level: int = 0,
    ) -> "ShapeElement":
        """
        Create a ShapeElement from a BaseShape.

        Args:
            slide_idx (int): The index of the slide.
            shape_idx (int): The index of the shape.
            shape (BaseShape): The shape object.
            config (Config): The configuration object.
            slide_area (float): The area of the slide.
            level (int): The indentation level.
            shape_cast (dict[MSO_SHAPE_TYPE, type[ShapeElement]] | None): Optional mapping of shape types to their corresponding ShapeElement classes.
            Set the value to None for any MSO_SHAPE_TYPE to exclude that shape type from processing.
        Returns:
            ShapeElement: The created ShapeElement.

        Raises:
            ValueError: If nested group shapes are not allowed.
        """
        if shape_idx > 100 and isinstance(shape, PPTXGroupShape):
            raise ValueError("Nested group shapes are not allowed")

        shape_normalize(shape)

        # Create style dictionary
        style = {
            "shape_bounds": {
                "width": shape.width,
                "height": shape.height,
                "left": shape.left,
                "top": shape.top,
            },
            "shape_type": str(shape.shape_type).split("(")[0].lower(),
            "rotation": shape.rotation,
            "name": shape.name,
        }

        # Determine semantic name
        try:
            # For auto shapes (rectangle, oval, triangle, star...)
            autoshape = shape.auto_shape_type
            assert autoshape is not None
            style["semantic_name"] = str(autoshape).split()[0].lower().strip()
        except Exception:
            # For other shapes (freeform, connector, table, chart...)
            style["semantic_name"] = str(shape.shape_type).split("(")[0].lower().strip()

        # Create text frame
        text_frame = TextFrame(shape, level + 1)

        # Create appropriate shape element based on shape type
        shape_class = shape_cast.get(shape.shape_type, UnsupportedShape)
        if shape_class is UnsupportedShape:
            shape_class = SHAPECAST.get(shape.shape_type, UnsupportedShape)

        if shape_class == Placeholder:
            shape_class = Placeholder.from_shape

        if shape_class == GroupShape:
            shape_class = GroupShape.with_shape_cast(shape_cast)

        return shape_class(
            config=config,
            slide_idx=slide_idx,
            shape_idx=shape_idx,
            style=style,
            data=[],
            text_frame=text_frame,
            level=level,
            slide_area=slide_area,
            xml=shape._element.xml,
            fill=Fill.from_shape(getattr(shape, "fill", None), shape.part, config),
            line=Line.from_shape(getattr(shape, "line", None), shape.part, config),
            shape=shape,
            _closures=ClosureType.to_default_dict(),
        )

    def build(self, slide: PPTXSlide) -> BaseShape:
        """
        Build the shape element in a slide.

        Args:
            slide (PPTXSlide): The slide to build the shape in.

        Returns:
            BaseShape: The built shape.
        """
        shape = slide.shapes._shape_factory(
            slide.shapes._spTree.insert_element_before(parse_xml(self.xml), "p:extLst")
        )
        if getattr(shape, "fill", None) is not None:
            self.fill.build(shape, shape.part)
        if getattr(shape, "line", None) is not None:
            self.line.build(shape.line, shape.part)
        return shape

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the shape element to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the shape element.

        Raises:
            NotImplementedError: If not implemented in a subclass.
        """
        raise NotImplementedError(
            f"to_html not implemented for {self.__class__.__name__}"
        )

    @property
    def text(self) -> str:
        """
        Get the text of the shape element.
        """
        if self.text_frame.is_textframe:
            return self.text_frame.text
        return ""

    def __getstate__(self) -> object:
        state = self.__dict__.copy()
        state["shape"] = None
        return state

    def __repr__(self) -> str:
        """
        Get a string representation of the shape element.

        Returns:
            str: A string representation of the shape element.
        """
        return f"{self.__class__.__name__}: shape {self.shape_idx} of slide {self.slide_idx}"

    @property
    def closures(self) -> list[Closure]:
        """
        Get the closures associated with the shape element.

        Returns:
            List[Closure]: A list of closures.
        """
        closures = []
        closures.extend(sorted(self._closures[ClosureType.CLONE]))
        closures.extend(
            self._closures[ClosureType.REPLACE] + self._closures[ClosureType.STYLE]
        )
        closures.extend(sorted(self._closures[ClosureType.DELETE], reverse=True))
        closures.extend(self._closures[ClosureType.MERGE])
        return closures

    @property
    def indent(self) -> str:
        """
        Get the indentation string for the shape element.

        Returns:
            str: The indentation string.
        """
        return "\t" * self.level

    @property
    def left(self) -> float:
        """
        Get the left position of the shape element.

        Returns:
            float: The left position in points.
        """
        return self.style["shape_bounds"]["left"].pt

    @left.setter
    def left(self, value: float) -> None:
        """
        Set the left position of the shape element.

        Args:
            value (float): The left position in points.
        """
        self.style["shape_bounds"]["left"] = value

    @property
    def top(self) -> float:
        """
        Get the top position of the shape element.

        Returns:
            float: The top position in points.
        """
        return self.style["shape_bounds"]["top"].pt

    @top.setter
    def top(self, value: float) -> None:
        """
        Set the top position of the shape element.

        Args:
            value (float): The top position in points.
        """
        self.style["shape_bounds"]["top"] = value

    @property
    def width(self) -> float:
        """
        Get the width of the shape element.

        Returns:
            float: The width in points.
        """
        return self.style["shape_bounds"]["width"].pt

    @width.setter
    def width(self, value: float) -> None:
        """
        Set the width of the shape element.

        Args:
            value (float): The width in points.
        """
        self.style["shape_bounds"]["width"] = value

    @property
    def height(self) -> float:
        """
        Get the height of the shape element.

        Returns:
            float: The height in points.
        """
        return self.style["shape_bounds"]["height"].pt

    @height.setter
    def height(self, value: float) -> None:
        """
        Set the height of the shape element.

        Args:
            value (float): The height in points.
        """
        self.style["shape_bounds"]["height"] = value

    @property
    def area(self) -> float:
        """
        Get the area of the shape element.

        Returns:
            float: The area in square points.
        """
        return self.width * self.height

    @property
    def semantic_name(self) -> Optional[str]:
        """
        Get the semantic name of the shape element.

        Returns:
            Optional[str]: The semantic name, or None if not set.
        """
        return self.style.get("semantic_name", None)

    @semantic_name.setter
    def semantic_name(self, value: str) -> None:
        """
        Set the semantic name of the shape element.

        Args:
            value (str): The semantic name.
        """
        self.style["semantic_name"] = value

    def get_inline_style(self, style_args: StyleArg) -> str:
        """
        Get the inline style for the shape element.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The inline style string.
        """
        id_str = f" id='{self.shape_idx}'" if style_args.element_id else ""
        data_attrs = []
        styles = []

        # Add data attributes
        if style_args.area:
            data_attrs.append(
                f"data-relative-area={self.area*100/self.slide_area:.2f}%;"
            )
        if style_args.show_name:
            data_attrs.append(f"data-shapeName='{self.style['name']}'")
        if style_args.show_semantic_name and self.semantic_name is not None:
            data_attrs.append(f"data-semanticName='{self.semantic_name}'")

        # Add style attributes
        if style_args.size:
            styles.append(f"width: {self.width}pt; height: {self.height}pt;")
        if style_args.geometry:
            styles.append(f"left: {self.left}pt; top: {self.top}pt;")
        if style_args.font_style and self.text_frame.is_textframe:
            font_style = self.text_frame.font.to_style()
            if font_style:
                styles.append(font_style)

        # Combine attributes
        if len(styles) != 0:
            id_str += " style='" + " ".join(styles) + "'"
        if len(data_attrs) != 0:
            id_str += " " + " ".join(data_attrs)

        return id_str


@dataclass
class UnsupportedShape(ShapeElement):
    def __post_init__(self) -> None:
        """
        Initialize an UnsupportedShape.

        Raises:
            ValueError: Always, as the shape is unsupported.
        """
        raise ValueError(f"Unsupported shape {self.shape.shape_type}")


class TextBox(ShapeElement):
    """
    A class to represent a text box shape.
    """

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the text box to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the text box.
        """
        content = self.text_frame.to_html(style_args)
        if not style_args.show_content:
            content = ""
        if not content and not style_args.show_empty:
            return ""
        return (
            f"{self.indent}<div{self.get_inline_style(style_args)}>\n"
            + content
            + f"\n{self.indent}</div>\n"
        )


@dataclass
class Picture(ShapeElement):
    """
    A class to represent a picture shape.
    """

    def __post_init__(self):
        """
        Create a Picture from a PPTXPicture.

        Returns:
            Picture: The created Picture.

        Raises:
            ValueError: If the image type is unsupported.
        """
        img_path = pjoin(
            self.config.IMAGE_DIR,
            f"{self.shape.image.sha1}.{self.shape.image.ext}",
        )
        img_path = parsing_image(self.shape.image, img_path)

        # Add image style information
        self.style["img_style"] = {
            "crop_bottom": self.shape.crop_bottom,
            "crop_top": self.shape.crop_top,
            "crop_left": self.shape.crop_left,
            "crop_right": self.shape.crop_right,
        }
        self.data.extend([img_path, self.shape.name, None])  # [img_path, name, caption]

    def build(self, slide: PPTXSlide) -> PPTXPicture:
        """
        Build the picture in a slide.

        Args:
            slide (PPTXSlide): The slide to build the picture in.

        Returns:
            PPTXPicture: The built picture.
        """
        # Add picture to slide
        if self.is_table:
            return slide.shapes.add_table(
                self.row, self.col, **self.style["shape_bounds"]
            )

        shape = slide.shapes.add_picture(
            self.img_path,
            **self.style["shape_bounds"],
        )

        # Set properties
        shape.name = self.style["name"]
        dict_to_object(self.style["img_style"], shape.image)

        # Apply shape bounds and rotation
        dict_to_object(self.style["shape_bounds"], shape)
        if hasattr(shape, "rotation"):
            shape.rotation = self.style["rotation"]

        return shape

    @property
    def is_table(self) -> bool:
        return self.style.get("is_table", False)

    @is_table.setter
    def is_table(self, value: bool) -> None:
        self.style["is_table"] = value

    @property
    def grid(self) -> tuple[int, int]:
        assert self.is_table, "The shape is not a table."
        return self.row, self.col

    @grid.setter
    def grid(self, value: tuple[int, int]) -> None:
        assert self.is_table, "The shape is not a table."
        self.row, self.col = value

    @property
    def img_path(self) -> str:
        """
        Get the image path.

        Returns:
            str: The image path.
        """
        return self.data[0]

    @img_path.setter
    def img_path(self, img_path: str) -> None:
        """
        Set the image path.

        Args:
            img_path (str): The image path.
        """
        self.data[0] = img_path

    @property
    def caption(self) -> Optional[str]:
        """
        Get the caption.

        Returns:
            Optional[str]: The caption, or None if not set.
        """
        return self.data[2]

    @caption.setter
    def caption(self, caption: str) -> None:
        """
        Set the caption.

        Args:
            caption (str): The caption.
        """
        self.data[2] = caption

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the picture to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the picture.

        Raises:
            ValueError: If the caption is not found.
        """
        if not style_args.show_image:
            return ""
        if self.caption is None:
            raise ValueError(
                f"Caption not found for picture {self.shape_idx} of slide {self.slide_idx}"
            )
        return (
            self.indent
            + f"<img {self.get_inline_style(style_args)} alt='{self.caption}'/>"
        )


@dataclass
class GroupShape(ShapeElement):
    """
    A class to represent a group shape.
    """

    shape_cast: ClassVar[dict[MSO_SHAPE_TYPE, type[ShapeElement]]] = {}

    @classmethod
    def with_shape_cast(cls, shape_cast: dict[MSO_SHAPE_TYPE, type[ShapeElement]]):
        """
        Dynamically create a subclass of GroupShape with an isolated shape_cast.
        """
        new_cls = type(f"{cls.__name__}_Isolated_{id(shape_cast)}", (cls,), {})
        new_cls.shape_cast = MappingProxyType(shape_cast)
        return new_cls

    def __post_init__(self) -> None:
        """
        Initialize a GroupShape.
        """
        # Create shape elements for each shape in the group
        self.data = [
            ShapeElement.from_shape(
                self.slide_idx,
                (self.shape_idx + 1) * 100 + i,
                sub_shape,
                self.config,
                self.slide_area,
                self.shape_cast,
                level=self.level + 1,
            )
            for i, sub_shape in enumerate(self.shape.shapes)
            if self.shape_cast.get(sub_shape.shape_type, -1) is not None
            and sub_shape.visible
        ]

        # Apply shape bounds to each shape in the group
        for idx, shape_bounds in enumerate(parse_groupshape(self.shape)):
            if not self.shape.shapes[idx].visible:
                continue
            if self.shape_cast.get(self.shape.shapes[idx].shape_type, -1) is None:
                continue
            self.data[idx].style["shape_bounds"] = shape_bounds

    def build(self, slide: PPTXSlide) -> PPTXSlide:
        """
        Build the group shape in a slide.

        Args:
            slide (PPTXSlide): The slide to build the group shape in.

        Returns:
            PPTXSlide: The slide with the built group shape.
        """
        for shape in self.data:
            shape.build(slide)
        return slide

    def shape_filter(
        self, shape_type: type["ShapeElement"], return_father: bool = False
    ):
        """
        Iterate over all shapes in the group.

        Yields:
            ShapeElement: Each shape in the group.
        """
        for shape in self.data:
            if isinstance(shape, shape_type):
                if return_father:
                    yield (self, shape)
                else:
                    yield shape

    @property
    def shapes(self):
        return self.data

    def __eq__(self, __value: object) -> bool:
        """
        Check if two group shapes are equal.

        Args:
            __value (object): The object to compare with.

        Returns:
            bool: True if the group shapes are equal, False otherwise.
        """
        if not isinstance(__value, GroupShape) or len(self.data) != len(__value.data):
            return False
        for shape1, shape2 in zip(self.data, __value.data):
            if isinstance(shape1, type(shape2)):
                return False
        return True

    def __repr__(self) -> str:
        """
        Get a string representation of the group shape.

        Returns:
            str: A string representation of the group shape.
        """
        return f"{self.__class__.__name__}: {self.data}"

    def __iter__(self):
        return iter(self.data)

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the group shape to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the group shape.
        """
        content = "\n".join([shape.to_html(style_args) for shape in self.data])
        if not style_args.show_content:
            content = ""
        return (
            self.indent
            + f"<div {self.get_inline_style(style_args)} data-group-label='{self.group_label}'>\n"
            + content
            + "\n"
            + self.indent
            + "</div>\n"
        )

    @property
    def group_label(self) -> str:
        """
        Get the group label.

        Returns:
            str: The group label.
        """
        return getattr(self, "_group_label", f"group_{self.shape_idx}")

    @group_label.setter
    def group_label(self, value: str) -> None:
        """
        Set the group label.

        Args:
            value (str): The group label.
        """
        self._group_label = value


class FreeShape(ShapeElement):
    """
    A class to represent a free shape.
    """

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the free shape to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the free shape.
        """
        content = self.text_frame.to_html(style_args)
        if not content and not style_args.show_empty:
            return ""
        return (
            f"{self.indent}<div {self.get_inline_style(style_args)}>"
            + f"\n{content}"
            + f"\n{self.indent}</div>"
        )


@dataclass
class SemanticPicture(Picture):
    """
    A class to represent a semantic picture (table, chart, etc.).
    """

    def __post_init__(self):
        shape_type = str(self.shape.shape_type).split()[0]
        self.style["img_style"] = {}
        self.data = [
            package_join("resource", "pic_placeholder.png"),
            self.shape.name,
            f"This is a picture of {shape_type}",
        ]
        self.semantic_name = shape_type


class Placeholder:
    """
    A class to represent a placeholder shape.
    """

    @classmethod
    def from_shape(
        cls,
        config: Config,
        slide_idx: int,
        shape_idx: int,
        shape: SlidePlaceholder,
        **kwargs,
    ) -> Union[Picture, TextBox]:
        """
        Create a Placeholder from a SlidePlaceholder.

        Returns:
            Union[Picture, TextBox]: The created shape element.

        Raises:
            ValueError: If the placeholder type is unsupported.
            AssertionError: If the placeholder has multiple types.
        """
        # Ensure placeholder has only one type
        assert (
            sum(
                [
                    shape.has_text_frame,
                    shape.has_chart,
                    shape.has_table,
                    isinstance(shape, PlaceholderPicture),
                ]
            )
            == 1
        ), "Placeholder should have only one type"

        # Create appropriate shape based on placeholder type
        if isinstance(shape, PlaceholderPicture):
            return Picture(
                config=config,
                slide_idx=slide_idx,
                shape_idx=shape_idx,
                shape=shape,
                **kwargs,
            )
        elif shape.has_text_frame:
            return TextBox(
                config=config,
                slide_idx=slide_idx,
                shape_idx=shape_idx,
                shape=shape,
                **kwargs,
            )
        else:
            raise ValueError(f"Unsupported placeholder {shape.placeholder_type}")


# Define shape type mapping
SHAPECAST = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: FreeShape,
    MSO_SHAPE_TYPE.LINE: FreeShape,
    MSO_SHAPE_TYPE.PICTURE: Picture,
    MSO_SHAPE_TYPE.PLACEHOLDER: Placeholder,
    MSO_SHAPE_TYPE.GROUP: GroupShape,
    MSO_SHAPE_TYPE.TEXT_BOX: TextBox,
    MSO_SHAPE_TYPE.MEDIA: SemanticPicture,
    MSO_SHAPE_TYPE.TABLE: SemanticPicture,
    MSO_SHAPE_TYPE.CHART: SemanticPicture,
    MSO_SHAPE_TYPE.LINKED_PICTURE: SemanticPicture,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: SemanticPicture,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: SemanticPicture,
    MSO_SHAPE_TYPE.DIAGRAM: SemanticPicture,
    MSO_SHAPE_TYPE.CANVAS: SemanticPicture,
    MSO_SHAPE_TYPE.INK: SemanticPicture,
    MSO_SHAPE_TYPE.IGX_GRAPHIC: SemanticPicture,
    MSO_SHAPE_TYPE.WEB_VIDEO: SemanticPicture,
}
