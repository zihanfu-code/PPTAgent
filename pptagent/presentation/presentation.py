import traceback
from collections.abc import Generator
from typing import Literal, Optional

from pptx import Presentation as load_prs
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape as PPTXGroupShape
from pptx.slide import Slide as PPTXSlide

from pptagent.utils import Config, get_logger, package_join

from .shapes import (
    Background,
    GroupShape,
    Paragraph,
    Picture,
    ShapeElement,
    StyleArg,
)

# Type variable for ShapeElement subclasses

logger = get_logger(__name__)


class SlidePage:
    """
    A class to represent a slide page in a presentation.
    """

    def __init__(
        self,
        shapes: list[ShapeElement],
        backgrounds: list[Background],
        slide_idx: int,
        real_idx: int,
        slide_notes: Optional[str],
        slide_layout_name: Optional[str],
        slide_title: Optional[str],
        slide_width: int,
        slide_height: int,
    ):
        """
        Initialize a SlidePage.

        Args:
            shapes (List[ShapeElement]): The shapes in the slide.
            backgrounds (List[Background]): The backgrounds of the slide.
            slide_idx (int): The index of the slide.
            real_idx (int): The real index of the slide.
            slide_notes (Optional[str]): The notes of the slide.
            slide_layout_name (Optional[str]): The layout name of the slide.
            slide_title (Optional[str]): The title of the slide.
            slide_width (int): The width of the slide.
            slide_height (int): The height of the slide.
        """
        self.shapes = shapes
        self.backgrounds = backgrounds
        self.slide_idx = slide_idx
        self.real_idx = real_idx
        self.slide_notes = slide_notes
        self.slide_layout_name = slide_layout_name
        self.slide_title = slide_title
        self.slide_width = slide_width
        self.slide_height = slide_height

        # Assign group labels to group shapes
        groups_shapes_labels = []
        for shape in self.shape_filter(GroupShape):
            for group_shape in groups_shapes_labels:
                if group_shape == shape:
                    shape.group_label = group_shape.group_label
                    continue
            groups_shapes_labels.append(shape)
            shape.group_label = f"group_{len(groups_shapes_labels)}"

    @classmethod
    def from_slide(
        cls,
        slide: PPTXSlide,
        slide_idx: int,
        real_idx: int,
        slide_width: int,
        slide_height: int,
        config: Config,
        shape_cast: dict[MSO_SHAPE_TYPE, type[ShapeElement] | None],
    ) -> "SlidePage":
        """
        Create a SlidePage from a PPTXSlide.

        Args:
            slide (PPTXSlide): The slide object.
            slide_idx (int): The index of the slide.
            real_idx (int): The real index of the slide.
            slide_width (int): The width of the slide.
            slide_height (int): The height of the slide.
            config (Config): The configuration object.
            shape_cast (dict[MSO_SHAPE_TYPE, type[ShapeElement] | None]): Mapping of shape types to their corresponding ShapeElement classes.
            Set the value to None for any MSO_SHAPE_TYPE to exclude that shape type from processing.
        Returns:
            SlidePage: The created SlidePage.
        """
        backgrounds = [Background.from_slide(slide, config)]
        shapes = []
        for i, shape in enumerate(slide.shapes):
            if not shape.visible:
                continue
            if shape_cast.get(shape.shape_type, -1) is None:
                continue
            shapes.append(
                ShapeElement.from_shape(
                    slide_idx, i, shape, config, slide_width * slide_height, shape_cast
                )
            )
        for i, s in enumerate(shapes):
            if isinstance(s, Picture) and s.area / s.slide_area > 0.95:
                backgrounds.append(shapes.pop(i))

        slide_layout_name = slide.slide_layout.name if slide.slide_layout else None
        slide_title = slide.shapes.title.text if slide.shapes.title else None
        slide_notes = (
            slide.notes_slide.notes_text_frame.text
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame
            else None
        )

        return cls(
            shapes,
            backgrounds,
            slide_idx,
            real_idx,
            slide_notes,
            slide_layout_name,
            slide_title,
            slide_width,
            slide_height,
        )

    def build(self, slide: PPTXSlide) -> PPTXSlide:
        """
        Build the slide page in a slide.

        Args:
            slide (PPTXSlide): The slide to build the slide page in.

        Returns:
            PPTXSlide: The built slide.
        """
        # Remove existing placeholders
        for ph in slide.placeholders:
            ph.element.getparent().remove(ph.element)

        # Build backgrounds, shapes and apply closures
        for shape in sorted(self.backgrounds + self.shapes, key=lambda x: x.shape_idx):
            build_shape = shape.build(slide)
            for closure in shape.closures:
                try:
                    closure.apply(build_shape)
                except Exception as e:
                    raise ValueError(f"Failed to apply closures to slides: {e}")
        return slide

    def iter_paragraphs(self) -> Generator[Paragraph, None, None]:
        for shape in self:  # this considered the group shapes
            if not shape.text_frame.is_textframe:
                continue
            for para in shape.text_frame.paragraphs:
                if para.idx != -1:
                    yield para

    def shape_filter(
        self,
        shape_type: type[ShapeElement],
        from_groupshape: bool = True,
        return_father: bool = False,
    ) -> (
        Generator[ShapeElement, None, None]
        | Generator[tuple["SlidePage", ShapeElement], None, None]
    ):
        """
        Filter shapes in the slide by type.

        Args:
            shape_type (Type[ShapeElement]): The type of shapes to filter.
            shapes (Optional[List[ShapeElement]]): The shapes to filter.

        Yields:
            ShapeElement: The filtered shapes.
        """
        for shape in self.shapes:
            if isinstance(shape, shape_type):
                if return_father:
                    yield (self, shape)
                else:
                    yield shape
            elif from_groupshape and isinstance(shape, GroupShape):
                yield from shape.shape_filter(shape_type, return_father)

    def get_content_type(self) -> Literal["text", "image"]:
        """
        Get the content type of the slide.

        Returns:
            Literal["text", "image"]: The content type of the slide.
        """
        if len(list(self.shape_filter(Picture))) == 0:
            return "text"
        return "image"

    def to_html(self, style_args: Optional[StyleArg] = None, **kwargs) -> str:
        """
        Represent the slide page in HTML.

        Args:
            style_args (Optional[StyleArg]): The style arguments for HTML conversion.
            **kwargs: Additional arguments.

        Returns:
            str: The HTML representation of the slide page.
        """
        if style_args is None:
            style_args = StyleArg(**kwargs)
        shapes_html = [shape.to_html(style_args) for shape in self.shapes]
        shapes_html = [html for html in shapes_html if html]
        return "".join(
            [
                "<!DOCTYPE html>\n<html>\n",
                (f"<title>{self.slide_title}</title>\n" if self.slide_title else ""),
                f'<body style="width:{self.slide_width}pt; height:{self.slide_height}pt;">\n',
                "\n".join(shapes_html),
                "</body>\n</html>\n",
            ]
        )

    def to_text(self, show_image: bool = False) -> str:
        """
        Represent the slide page in text.

        Args:
            show_image (bool): Whether to show image captions.

        Returns:
            str: The text representation of the slide page.

        Raises:
            ValueError: If an image caption is not found.
        """
        text_content = ""
        for para in self.iter_paragraphs():
            if not para.text:
                continue
            if para.bullet:
                text_content += para.bullet
            text_content += para.text + "\n"
        if show_image:
            for image in self.shape_filter(Picture):
                text_content += "\n" + "Image: " + image.caption
        return text_content

    def __iter__(self):
        """
        Iterate over all shapes in the slide page.

        Yields:
            ShapeElement: Each shape in the slide page.
        """
        for shape in self.shapes:
            if isinstance(shape, GroupShape):
                yield from shape
            else:
                yield shape

    def __len__(self) -> int:
        """
        Get the number of shapes in the slide page.

        Returns:
            int: The number of shapes.
        """
        return len(self.shapes)


class Presentation:
    """
    PPTAgent's representation of a presentation.
    Aiming at a more readable and editable interface.
    """

    def __init__(
        self,
        slides: list[SlidePage],
        error_history: list[tuple[int, str]],
        slide_width: float,
        slide_height: float,
        file_path: str,
        num_pages: int,
    ) -> None:
        """
        Initialize the Presentation.

        Args:
            slides (List[SlidePage]): The slides in the presentation.
            error_history (List[Tuple[int, str]]): The error history.
            slide_width (float): The width of the slides.
            slide_height (float): The height of the slides.
            file_path (str): The path to the presentation file.
            num_pages (int): The number of pages in the presentation.
        """
        self.slides = slides
        self.error_history = error_history
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.num_pages = num_pages
        self.source_file = file_path
        self.prs = load_prs(self.source_file)
        self.layout_mapping = {layout.name: layout for layout in self.prs.slide_layouts}
        self.prs.core_properties.last_modified_by = "PPTAgent"

    @classmethod
    def from_file(
        cls,
        file_path: str,
        config: Config,
        shape_cast: Optional[dict[MSO_SHAPE_TYPE, type[ShapeElement]]] = None,
    ) -> "Presentation":
        """
        Parse a Presentation from a file.

        Args:
            file_path (str): The path to the presentation file.
            config (Config): The configuration object.
            shape_cast (dict[MSO_SHAPE_TYPE, type[ShapeElement]] | None): Optional mapping of shape types to their corresponding ShapeElement classes.
            Set the value to None for any MSO_SHAPE_TYPE to exclude that shape type from processing.
        Returns:
            Presentation: The parsed Presentation.
        """
        prs = load_prs(file_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slides = []
        error_history = []
        slide_idx = 0
        layouts = [layout.name for layout in prs.slide_layouts]
        num_pages = len(prs.slides)

        if shape_cast is None:
            shape_cast = {}

        for slide in prs.slides:
            # Skip slides that won't be printed to PDF, as they are invisible
            if slide._element.get("show", 1) == "0":
                continue

            slide_idx += 1
            try:
                if slide.slide_layout.name not in layouts:
                    raise ValueError(
                        f"Slide layout {slide.slide_layout.name} not found"
                    )
                slides.append(
                    SlidePage.from_slide(
                        slide,
                        slide_idx - len(error_history),
                        slide_idx,
                        slide_width.pt,
                        slide_height.pt,
                        config,
                        shape_cast,
                    )
                )
            except Exception as e:
                error_history.append((slide_idx, str(e)))
                logger.warning(
                    "Fail to parse slide %d of %s: %s",
                    slide_idx,
                    file_path,
                    e,
                )
                logger.warning(traceback.format_exc())

        return cls(
            slides, error_history, slide_width, slide_height, file_path, num_pages
        )

    def save(self, file_path: str, layout_only: bool = False) -> None:
        """
        Save the presentation to a file.

        Args:
            file_path (str): The path to save the presentation to.
            layout_only (bool): Whether to save only the layout.
        """
        self.clear_slides()
        for slide in self.slides:
            if layout_only:
                self.clear_images(slide.shapes)
            pptx_slide = self.build_slide(slide)
            if layout_only:
                self.clear_text(pptx_slide.shapes)
        self.prs.save(file_path)

    def build_slide(self, slide: SlidePage) -> PPTXSlide:
        """
        Build a slide in the presentation.
        """
        return slide.build(
            self.prs.slides.add_slide(self.layout_mapping[slide.slide_layout_name])
        )

    def clear_slides(self):
        """
        Delete all slides from the presentation.
        """
        while len(self.prs.slides) != 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def clear_images(self, shapes: list[ShapeElement]):
        for shape in shapes:
            if isinstance(shape, GroupShape):
                self.clear_images(shape.shapes)
            elif isinstance(shape, Picture):
                shape.img_path = package_join("resource", "pic_placeholder.png")

    def clear_text(self, shapes: list[BaseShape]):
        for shape in shapes:
            if isinstance(shape, PPTXGroupShape):
                self.clear_text(shape.shapes)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "a" * len(run.text)

    def to_text(self, show_image: bool = False) -> str:
        """
        Represent the presentation in text.
        """
        return "\n----\n".join(
            [
                (
                    f"Slide {slide.slide_idx} of {len(self.slides)}\n"
                    + (f"Title:{slide.slide_title}\n" if slide.slide_title else "")
                    + slide.to_text(show_image)
                )
                for slide in self.slides
            ]
        )

    def __iter__(self):
        yield from self.slides

    def __len__(self) -> int:
        """
        Get the number of slides in the presentation.
        """
        return len(self.slides)

    def __getstate__(self) -> object:
        state = self.__dict__.copy()
        state["prs"] = None
        state["layout_mapping"] = None
        return state

    def __setstate__(self, state: object):
        self.__dict__.update(state)
        self.prs = load_prs(self.source_file)
        self.layout_mapping = {layout.name: layout for layout in self.prs.slide_layouts}
