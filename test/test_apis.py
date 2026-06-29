from test.conftest import test_config

from bs4 import BeautifulSoup
from pptx import Presentation

from pptagent.apis import (
    API_TYPES,
    CodeExecutor,
    markdown,
    process_element,
    replace_para,
)


def test_api_docs():
    executor = CodeExecutor(3)
    docs = executor.get_apis_docs(API_TYPES.Agent.value)
    assert len(docs) > 0


def test_replace_para():
    text = "这是一个**加粗和*斜体*文本**，还有*斜体和`Code def a+b`*，~~删除~~，[链接](http://example.com)"
    prs = Presentation(test_config.ppt)
    slide = prs.slides[0]
    replace_para(0, text, slide.shapes[0])
    runs = slide.shapes[0].text_frame.paragraphs[0].runs
    assert runs[1].font.bold
    assert runs[2].font.bold and runs[2].font.italic
    assert runs[6].font.name == "Consolas"
    assert runs[8].font.strikethrough
    assert runs[10].hyperlink.address == "http://example.com"


def test_list_parsing():
    text = """
    - 项目1
    - 项目2

    1. 项目1
    2. 项目2
    """
    html = markdown(text).strip()
    soup = BeautifulSoup(html, "html.parser")
    blocks = process_element(soup)
    assert len(blocks) == 1
    assert "ol" not in html and "ul" not in html
